"""
Document loaders.

Every loader returns a list of "page records":
    {"text": str, "page": int, "source": str, "metadata": {...}}

Keeping page-level granularity (rather than one giant string per file) lets
the chunker preserve accurate page citations later, which is critical for
the "source citation" requirement of the assistant.

Supported formats: .pdf, .docx, .pptx, .txt, .md
Each loader also tries to parse a small metadata header if present, in the
form:
    TITLE: ...
    DEPARTMENT: ...
    SENSITIVITY: ...
    VERSION: ...
    EFFECTIVE_DATE: ...
This mimics how a real enterprise pipeline would pull metadata either from
document properties or from a companion metadata.json / SharePoint column.
Falls back to filename + file system info when absent.
"""
from __future__ import annotations

import os
import re
from datetime import datetime
from pathlib import Path
from typing import Any


HEADER_FIELDS = ["TITLE", "DEPARTMENT", "SENSITIVITY", "VERSION", "EFFECTIVE_DATE"]


def _parse_header_metadata(text: str) -> dict[str, str]:
    """Pull TITLE/DEPARTMENT/SENSITIVITY/VERSION/EFFECTIVE_DATE lines from the
    top of a document if the author included them. Returns {} if none found."""
    meta = {}
    for line in text.splitlines()[:15]:
        for field in HEADER_FIELDS:
            if line.strip().upper().startswith(field + ":"):
                meta[field.lower()] = line.split(":", 1)[1].strip()
    return meta


def _base_metadata(filepath: str) -> dict[str, Any]:
    p = Path(filepath)
    stat = p.stat()
    return {
        "source": p.name,
        "filepath": str(p),
        "file_type": p.suffix.lower().lstrip("."),
        "ingested_at": datetime.utcnow().isoformat(),
        "last_modified": datetime.utcfromtimestamp(stat.st_mtime).isoformat(),
        "organization": "Northwind Analytics, Inc.",   # <-- ADD THIS LINE (default/parent org)
        "department": "General",
        "sensitivity": "Internal",
        "title": p.stem,
        "version": "1.0",
    }


def load_txt(filepath: str) -> list[dict[str, Any]]:
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    meta = _base_metadata(filepath)
    meta.update(_parse_header_metadata(text))
    # .txt/.md have no real "pages" -- treat whole file as page 1
    return [{"text": text, "page": 1, "source": meta["source"], "metadata": meta}]


def load_pdf(filepath: str) -> list[dict[str, Any]]:
    from pypdf import PdfReader

    reader = PdfReader(filepath)
    meta = _base_metadata(filepath)
    # try metadata header from first page text
    first_page_text = reader.pages[0].extract_text() or "" if reader.pages else ""
    meta.update(_parse_header_metadata(first_page_text))

    records = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            records.append({"text": text, "page": i, "source": meta["source"], "metadata": meta})
    return records


def load_docx(filepath: str) -> list[dict[str, Any]]:
    import docx  # python-docx

    doc = docx.Document(filepath)
    full_text = "\n".join(p.text for p in doc.paragraphs)
    meta = _base_metadata(filepath)
    meta.update(_parse_header_metadata(full_text))

    # Word has no reliable page boundaries without rendering; we chunk the
    # whole doc as one logical "page" and let the chunker split it further.
    # For tables, append a simple pipe-delimited representation.
    tables_text = []
    for t_idx, table in enumerate(doc.tables, start=1):
        rows = [" | ".join(cell.text for cell in row.cells) for row in table.rows]
        tables_text.append(f"[Table {t_idx}]\n" + "\n".join(rows))
    if tables_text:
        full_text += "\n\n" + "\n\n".join(tables_text)

    return [{"text": full_text, "page": 1, "source": meta["source"], "metadata": meta}]


def load_pptx(filepath: str) -> list[dict[str, Any]]:
    from pptx import Presentation

    prs = Presentation(filepath)
    meta = _base_metadata(filepath)

    records = []
    all_text_for_header = ""
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    texts.append(" | ".join(c.text for c in row.cells))
        slide_text = "\n".join(t for t in texts if t.strip())
        if i == 1:
            all_text_for_header = slide_text
        if slide_text.strip():
            records.append({"text": slide_text, "page": i, "source": meta["source"], "metadata": meta})

    meta.update(_parse_header_metadata(all_text_for_header))
    for r in records:
        r["metadata"] = meta
    return records


LOADER_MAP = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".txt": load_txt,
    ".md": load_txt,
}


def load_document(filepath: str) -> list[dict[str, Any]]:
    ext = Path(filepath).suffix.lower()
    if ext not in LOADER_MAP:
        raise ValueError(f"Unsupported file type: {ext} ({filepath})")
    return LOADER_MAP[ext](filepath)


def load_directory(directory: str) -> list[dict[str, Any]]:
    """Walk a directory and load every supported file. Skips unsupported /
    unreadable files with a warning instead of crashing the whole ingest run
    -- important for a real enterprise repo with mixed junk files."""
    records = []
    for root, _, files in os.walk(directory):
        for fname in files:
            ext = Path(fname).suffix.lower()
            if ext not in LOADER_MAP:
                continue
            fpath = os.path.join(root, fname)
            try:
                records.extend(load_document(fpath))
            except Exception as e:
                print(f"[WARN] Failed to load {fpath}: {e}")
    return records
